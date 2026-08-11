from pptx import Presentation

PPTX_PATH = r"C:\Temp\Github Copilot and Nigel.pptx"
OUT_PATH = PPTX_PATH
FALLBACK_OUT = r"C:\Temp\Github Copilot and Nigel_polished.pptx"

# Exact run-text replacements only (preserves shape/layout/pictures and run formatting)
REPLACEMENTS = {
    "Different models have different \u201cdiscount\u201c. ": "Different models have different cost multipliers.",
    "Calling time consuming method ": "Calling the time-consuming method ",
    "Github": "GitHub",
    " Check the updated source code": "Reviewed the updated source code",
    " changed from 60 minutes to 10 minutes. ": " was reduced from 60 minutes to 10 minutes.",
    " Bonus, these slides are also helped by ": "Bonus: these slides were also created with ",
    "RF Test Results Don't Always Match Expectation": "RF Test Results Do Not Always Match Expectations",
    " and asks different people for help": " and asking different people for help",
    "topic-playbook.md    \u2014 domain debug guidance (EVM, Gain, PAE, IP3, SEM, DPD, ...)": "topic-playbook.md    \u2014 domain troubleshooting guidance (EVM, Gain, PAE, IP3, SEM, DPD, ...)",
    " \u201c/rf-troubleshooting\u201d to reference this skill, then Copilot will trouble shooting step-by-step with user input information": "Use \"/rf-troubleshooting\" to invoke this skill, then Copilot troubleshoots step by step based on user inputs",
    "Skills stored in the markdown file can be incrementally enhanced by iteratively adding new expert knowledge": "Skills stored in Markdown files can be incrementally improved by continuously adding expert knowledge",
    "The \u201cskills\u201d are maintained as human-readable *.md files and can be generated automatically using AI.": "Skills are maintained as human-readable *.md files and can also be generated automatically with AI.",
    "Very convenient to compare different AI models with the same prompt and the same skills": "It is very convenient to compare different AI models using the same prompt and the same skills",
    "Context-aware answers uses LabVIEW idioms, ": "Context-aware answers use LabVIEW idioms, ",
    "Instant summaries - explains APIs, functions, or step types without digging through manuals": "Instant summaries explain APIs, functions, or step types without digging through manuals",
    "Make Review Code Faster": "Make Code Reviews Faster",
    "Large VIs,": "Large VIs, ",
    "c": "",  # only safe in this deck for the known split word before 'ustom '
    "ustom ": "custom ",
    "Identifying - dead code, maintainability issues, performance risks": "Identifying dead code, maintainability issues, and performance risks",
    "Code explanation, summarizes what a VI or sequence is doing, step by step": "Code explanations summarize what a VI or sequence is doing, step by step",
    "Pattern recognition, identifies common architectures (State Machine, Producer/Consumer, Actor Framework)": "Pattern recognition identifies common architectures (State Machine, Producer/Consumer, Actor Framework)",
    " Develop Faster with LabVIEW and ": "Develop Faster with LabVIEW and ",
    "Rewriting boilerplate code like error in error out wiring in LabVIEW, standard sequence structure like Setup, Main, Cleanup in ": "Rewriting boilerplate code such as error-in/error-out wiring in LabVIEW and standard sequence structures like Setup, Main, and Cleanup in ",
    "Code generation, use skeletons for": "Code generation uses skeletons for",
    " Debugging support": "Debugging support",
    "\u2705 Faster peer reviews": "\u2705 Faster implementation cycles",
    "\u2705 Easier handover of legacy projects": "\u2705 More consistent architecture decisions",
    "\u2705 Higher code quality with less effort": "\u2705 Higher code quality with less rework",
    "Example \u2013 Uses Nigel in LabVIEW 2026": "Example \u2013 Using Nigel in LabVIEW 2026",
    "Example \u2013 Uses Nigel in ": "Example \u2013 Using Nigel in ",
}

# Context-specific replacements to avoid unsafe global changes for short tokens
CONTEXT_REPLACEMENTS = [
    # Slide 4 line polishing
    (4, "The most recent models are Claude 4.6 and GPT 5.4", "Recent model options include Claude 4.6 and GPT-5.4"),
    # Slide 6 goal sentence
    (6, " Using ", "Using "),
    (6, " Copilot to refactor DPAT version to handle large files efficiently while preserving all features", " Copilot to refactor the DPAT version for efficient large-file handling while preserving all features"),
    # Slide 7 key benefit sentence fragments
    (7, "stdf", "STDF"),
    (7, "real project ", "production project "),
    # Slide 10 sentence fragment
    (10, "take more time", "spend more time"),
]

prs = Presentation(PPTX_PATH)
changes = 0

for s_idx, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                original = run.text
                text = original

                # Exact replacements
                for old, new in REPLACEMENTS.items():
                    if text == old:
                        text = new

                # Context replacements by slide index
                for target_slide, old, new in CONTEXT_REPLACEMENTS:
                    if s_idx == target_slide and text == old:
                        text = new

                if text != original:
                    run.text = text
                    changes += 1

try:
    prs.save(OUT_PATH)
    print(f"Saved in place: {OUT_PATH}")
except PermissionError:
    prs.save(FALLBACK_OUT)
    print(f"Original file is locked. Saved to: {FALLBACK_OUT}")

print(f"Run-level text updates applied: {changes}")
