"""
Inserts 4 new slides after slide 9 in "Github Copilot and Nigel.pptx".

Slide order inserted (at index 9, pushing current 10-16 to 14-20):
  A. The Challenge: RF Test Troubleshooting Takes Time        (problem intro)
  B. Step 1: Building the RF Troubleshooting Skill           (skill creation)
  C. Step 2: Using the Skill — WLAN EVM Example              (usage demo)
  D. Step 3: How GitHub Copilot + Skill Helped               (value summary)
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import lxml.etree as etree

PPTX_PATH = r"C:\Repos\Semi-Skills-Hub\skills\rf-troubleshooting\Github Copilot and Nigel.pptx"

prs = Presentation(PPTX_PATH)

# ── Layout to use (matches slides 9 and 10) ────────────────────────────────
LAYOUT = prs.slide_layouts[12]   # "1_Plain Content Slide"

# ── Color constants (NI / deck theme) ──────────────────────────────────────
BLUE    = RGBColor(0x00, 0x70, 0xC0)   # section header colour
DARK    = RGBColor(0x26, 0x26, 0x26)   # body text
GREEN   = RGBColor(0x70, 0xAD, 0x47)   # positive/feature highlight
RED_STR = RGBColor(0xC0, 0x00, 0x00)   # problem/negative highlight
GRAY    = RGBColor(0x80, 0x80, 0x80)   # secondary


def set_font(run, size_pt=None, bold=None, color=None):
    if size_pt:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_heading(tf, text, size_pt=20, color=BLUE, space_before_pt=10):
    """Add a section-header paragraph."""
    para = tf.add_paragraph()
    para.space_before = Pt(space_before_pt)
    para.level = 0
    run = para.add_run()
    run.text = text
    set_font(run, size_pt=size_pt, bold=True, color=color)
    return para


def add_bullet(tf, text, level=1, size_pt=16, bold=False, color=DARK):
    """Add a bullet paragraph."""
    para = tf.add_paragraph()
    para.level = level
    run = para.add_run()
    run.text = text
    set_font(run, size_pt=size_pt, bold=bold, color=color)
    return para


def add_blank(tf):
    tf.add_paragraph()


def create_slide(title_text):
    """Create a new slide using the standard layout and set the title."""
    slide = prs.slides.add_slide(LAYOUT)
    # Set title
    for shape in slide.shapes:
        if shape.has_text_frame and "Title" in shape.name:
            shape.text = title_text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
            break
    # Clear the content placeholder
    for shape in slide.shapes:
        if "Content" in shape.name and shape.has_text_frame:
            # Remove all existing paragraphs except one
            tf = shape.text_frame
            tf.clear()
            return slide, tf
    return slide, None


def move_slide(prs, old_index, new_index):
    """Move a slide from old_index to new_index in the presentation XML."""
    xml_slides = prs.slides._sldIdLst
    slides     = list(xml_slides)
    elem       = slides[old_index]
    xml_slides.remove(elem)
    xml_slides.insert(new_index, elem)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE A — The Challenge
# ═══════════════════════════════════════════════════════════════════════════
slideA, tf = create_slide("The Challenge: RF Test Troubleshooting Takes Time")

add_heading(tf, "RF Test Results Don't Always Match Expectation", size_pt=20, space_before_pt=0)
add_bullet(tf, "Wide variety of root causes: measurement config, calibration, waveform, timing, or DUT", level=1)
add_bullet(tf, "Engineers unfamiliar with this area may spend hours checking irrelevant paths", level=1, bold=True, color=RED_STR)
add_bullet(tf, "Proper isolation requires deep domain knowledge across 10+ RF test types", level=1)

add_blank(tf)
add_heading(tf, "Example: WLAN EVM Measured 3 dB Worse Than Expected", size_pt=20)
add_bullet(tf, "Possible causes: offset misalignment, measurement length, channel estimation, reference level, waveform...", level=1)
add_bullet(tf, "Without guidance: where do you start?  →  Senior engineer? PDF manual? Trial and error?", level=1, bold=True, color=RED_STR)

add_blank(tf)
add_heading(tf, "The Cost", size_pt=20)
add_bullet(tf, "Traditional approach: search ~1,000 pages of PDFs, Confluence, ask senior engineers", level=1)
add_bullet(tf, "Slow, inconsistent, expertise bottlenecked to a few individuals", level=1)
add_bullet(tf, "GitHub Copilot alone gives generic guesses — useful for code, but not for structured RF fault isolation", level=1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE B — Building the Skill
# ═══════════════════════════════════════════════════════════════════════════
slideB, tf = create_slide("Step 1: Building the RF Troubleshooting Skill")

add_heading(tf, "📄  Input: Raw Expert Documents", size_pt=20, space_before_pt=0)
add_bullet(tf, "Reference Solution for Semi RFIC Production Test  (PDF)", level=1)
add_bullet(tf, "RF FEM APT Test Manual  (PDF)", level=1)
add_bullet(tf, "~1,000 pages of NI RF production-test knowledge", level=1, color=GRAY)

add_blank(tf)
add_heading(tf, "🔧  Process: Distill into 5 Structured Reference Files", size_pt=20)
add_bullet(tf, "coverage-map.md      — maps PDF chapters to reference files", level=1, color=DARK)
add_bullet(tf, "topic-playbook.md    — domain debug guidance (EVM, Gain, PAE, IP3, SEM, DPD, ...)", level=1, color=DARK)
add_bullet(tf, "symptom-taxonomy.md  — symptom → first checks", level=1, color=DARK)
add_bullet(tf, "isolation-workflow.md — step-by-step fault isolation", level=1, color=DARK)
add_bullet(tf, "troubleshooting-summary.md — cross-cutting principles and heuristics", level=1, color=DARK)

add_blank(tf)
add_heading(tf, "🚀  Deploy: Copy Skill to VS Code", size_pt=20)
add_bullet(tf, "Personal:  ~/.copilot/skills/rf-troubleshooting/", level=1)
add_bullet(tf, "Project:   .github/skills/rf-troubleshooting/  (scoped to one repo)", level=1)
add_bullet(tf, "Invoke:    Attach SKILL.md in Copilot Chat as  #prompt:SKILL.md", level=1, bold=True, color=BLUE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE C — Using the Skill
# ═══════════════════════════════════════════════════════════════════════════
slideC, tf = create_slide("Step 2: Using the Skill — WLAN EVM Example")

add_heading(tf, "Problem Statement", size_pt=20, space_before_pt=0)
add_bullet(tf, "WLAN DUT measured EVM = −38 dB,  expected = −35 dB  (3 dB worse than expected)", level=1, bold=True, color=RED_STR)

add_blank(tf)
add_heading(tf, "① Invoke — Attach SKILL.md in Copilot Chat", size_pt=18)
add_bullet(tf, '#prompt:SKILL.md  →  "WLAN EVM 3 dB worse than expected, how to troubleshoot?"', level=1, color=DARK)

add_heading(tf, "② Missing-Info Gate — Copilot Asks First (Structured Dialogs)", size_pt=18)
add_bullet(tf, "Expected vs measured EVM?  →  −35 dB expected, −38 dB measured", level=1)
add_bullet(tf, "WLAN standard?  →  802.11b/g", level=1)
add_bullet(tf, "Measurement offset configured?  →  Yes, symbol-based offset", level=1)
add_bullet(tf, "Loopback result?  →  Also 3 dB worse  ← KEY PIVOT", level=1, bold=True, color=BLUE)

add_heading(tf, "③ Loopback Worse = Path / Cable / DUT Eliminated Immediately", size_pt=18, color=GREEN)
add_bullet(tf, "Focus shifts to: measurement configuration layer only", level=1, bold=True, color=GREEN)

add_heading(tf, "④ Copilot Delivers 5 Prioritized, Actionable Checks", size_pt=18)
add_bullet(tf, "1. Verify SymbolOffset alignment (use RFmx Soft Front Panel)", level=1)
add_bullet(tf, "2. Increase MeasurementLength to ≥10 symbols", level=1)
add_bullet(tf, "3. Compare ChannelEstimationMode vs bench reference", level=1)
add_bullet(tf, "4. Verify ReferenceLevel is 5 dB above average signal power", level=1)
add_bullet(tf, "5. Re-export bench waveform, reload on ATE, re-run loopback", level=1)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE D — How It Helped
# ═══════════════════════════════════════════════════════════════════════════
slideD, tf = create_slide("Step 3: How GitHub Copilot + Skill Helped")

add_heading(tf, "❌  Without the Skill", size_pt=20, color=RED_STR, space_before_pt=0)
add_bullet(tf, "Generic Copilot: 15+ vague suggestions (noise, hardware, DUT, drivers...)", level=1, color=RED_STR)
add_bullet(tf, "No prioritization, no production-test context, no structured isolation path", level=1, color=RED_STR)
add_bullet(tf, "Engineer spends hours verifying irrelevant paths before reaching root cause", level=1, color=RED_STR)

add_blank(tf)
add_heading(tf, "✅  With the Skill", size_pt=20, color=GREEN)
add_bullet(tf, "Asked for loopback result first → single most discriminating fact collected immediately", level=1, color=GREEN)
add_bullet(tf, "Loopback = wrong → eliminated path, cable, calibration, and DUT as root causes in one step", level=1, color=GREEN)
add_bullet(tf, "3 prioritized measurement-config checks delivered, root cause isolated in one session", level=1, color=GREEN, bold=True)

add_blank(tf)
add_heading(tf, "Why It Works", size_pt=20)
add_bullet(tf, "Missing-Info Gate:    asked loopback first → eliminated entire failure layers", level=1)
add_bullet(tf, "Symptom Taxonomy:  'loopback also wrong' → mapped to measurement-config layer", level=1)
add_bullet(tf, "Topic Playbook:        802.11b/g-specific checks: offset, length, channel estimation", level=1)
add_bullet(tf, "Source-backed:        RF FEM APT Manual p.78–81 · Reference Solution p.10–14", level=1, color=GRAY)
add_bullet(tf, "Covers all 18 test domains: Gain, PAE, EVM, IP3, SEM, DPD, TTR, DPAT, and more", level=1, bold=True, color=BLUE)

# ═══════════════════════════════════════════════════════════════════════════
# INSERT all 4 slides at position 9 (after slide 9, before old slide 10)
# They were added at the end (indices 16, 17, 18, 19).
# After moving each to index 9, the final order will be reversed, so move in
# reverse order: D → C → B → A to end up with A-B-C-D at 9-12.
# ═══════════════════════════════════════════════════════════════════════════
total = len(prs.slides.element.findall(
    ".//{http://schemas.openxmlformats.org/presentationml/2006/main}sldId"))

move_slide(prs, total - 1, 9)  # D  (was last, now at 9; old 9→10, 10→11, ...)
move_slide(prs, total - 1, 9)  # C  (was last-1, now at 9; D shifts to 10)
move_slide(prs, total - 1, 9)  # B  (was last-2, now at 9; C→10, D→11)
move_slide(prs, total - 1, 9)  # A  (was last-3, now at 9; B→10, C→11, D→12)

# ── Verify final slide order ───────────────────────────────────────────────
print("Final slide order:")
for i, s in enumerate(prs.slides):
    title = ""
    for shape in s.shapes:
        if shape.has_text_frame and "Title" in shape.name:
            title = shape.text.strip()[:70]
            break
    print(f"  [{i+1:02d}] {title}")

OUT_PATH = PPTX_PATH.replace(".pptx", "_updated.pptx")
prs.save(OUT_PATH)
print("\nSaved:", OUT_PATH)
